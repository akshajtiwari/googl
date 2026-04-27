from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

OUT_FILE = "Solution_Challenge_2026_Prototype_Auto.pptx"

# Theme from lib/core/theme.dart
PRIMARY = RGBColor(0x1A, 0x3A, 0x5C)
ACCENT = RGBColor(0x2E, 0x5F, 0x8A)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF9)
TEXT = RGBColor(0x2D, 0x37, 0x48)
MUTED = RGBColor(0x71, 0x81, 0x96)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
YELLOW = RGBColor(0xF5, 0x9E, 0x0B)
TEAL = RGBColor(0x06, 0xB6, 0xD4)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_banner(slide, title, subtitle=""):
    # Persistent upper banner as requested
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.72))
    banner.fill.solid()
    banner.fill.fore_color.rgb = PRIMARY
    banner.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.69), prs.slide_width, Inches(0.03))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.35), Inches(0.10), Inches(8.8), Inches(0.45))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = WHITE

    if subtitle:
        st = slide.shapes.add_textbox(Inches(9.0), Inches(0.14), Inches(4.0), Inches(0.35))
        stf = st.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.RIGHT
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.size = Pt(11)
        sr.font.color.rgb = RGBColor(0xD1, 0xDE, 0xEA)


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(12.2), Inches(0.8))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    if subtitle:
        sp = tf.add_paragraph()
        sp.text = subtitle
        sp.font.size = Pt(13)
        sp.font.color.rgb = MUTED


def add_bullets(slide, x, y, w, h, items, title=None):
    if title:
        t = slide.shapes.add_textbox(x, y - Inches(0.25), w, Inches(0.25))
        tp = t.text_frame.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(16)
        tp.font.bold = True
        tp.font.color.rgb = PRIMARY

    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.clear()
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT


def add_kpi_card(slide, x, y, w, h, label, value, color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(1.75)

    vb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.12), w - Inches(0.3), Inches(0.45))
    vp = vb.text_frame.paragraphs[0]
    vp.text = value
    vp.font.size = Pt(24)
    vp.font.bold = True
    vp.font.color.rgb = color

    lb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.62), w - Inches(0.3), Inches(0.35))
    lp = lb.text_frame.paragraphs[0]
    lp.text = label
    lp.font.size = Pt(12)
    lp.font.color.rgb = MUTED


def add_flow_box(slide, x, y, w, h, text, color=ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = color
    s.line.width = Pt(1.5)
    tf = s.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    return s


def add_arrow(slide, x, y, w=Inches(0.45), h=Inches(0.18)):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid()
    a.fill.fore_color.rgb = ACCENT
    a.line.fill.background()


# Slide 1: Cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_banner(slide, "Solution Challenge 2026 Prototype", "Upper banner preserved")
add_title(slide, "NGO Coordinator Platform", "AI-assisted dashboard for needs mapping, assignment, and impact analytics")
add_bullets(
    slide,
    Inches(0.6), Inches(2.1), Inches(8.2), Inches(2.6),
    [
        "Project source: Flutter admin portal (`googl/`) with dashboard, heatmap, tasks, assignments, and analytics modules.",
        "Purpose: reduce response time, increase volunteer-task match quality, and improve zone-level coverage.",
        "Focus area demonstrated in this prototype: emergency community support operations in Delhi zones.",
    ],
)
add_kpi_card(slide, Inches(9.2), Inches(2.1), Inches(3.7), Inches(1.2), "Total Volunteers", "142", BLUE)
add_kpi_card(slide, Inches(9.2), Inches(3.45), Inches(3.7), Inches(1.2), "Active Tasks", "38", GREEN)
add_kpi_card(slide, Inches(9.2), Inches(4.8), Inches(3.7), Inches(1.2), "Zone Coverage", "74%", TEAL)

# Slide 2: Problem
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_banner(slide, "Problem Statement", "Need-response coordination gaps")
add_title(slide, "Why Existing Coordination Breaks Under Load")
add_bullets(
    slide,
    Inches(0.6), Inches(1.8), Inches(6.2), Inches(4.8),
    [
        "Need reports are fragmented across zones and channels.",
        "Manual volunteer assignment leads to skill and proximity mismatch.",
        "Underserved zones remain unnoticed until escalations happen.",
        "Coordinators lack a single place for timeline, status, and verification.",
        "Stakeholder reporting is time-consuming without live analytics.",
    ],
    title="Observed Pain Points"
)
add_flow_box(slide, Inches(7.1), Inches(2.0), Inches(5.6), Inches(0.8), "21 pending needs + 38 active tasks -> coordination overhead", RED)
add_flow_box(slide, Inches(7.1), Inches(3.1), Inches(5.6), Inches(0.8), "Overdue task risk when task ownership is unclear", YELLOW)
add_flow_box(slide, Inches(7.1), Inches(4.2), Inches(5.6), Inches(0.8), "Coverage imbalance across wards without heatmap overlays", BLUE)

# Slide 3: Solution Overview
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_banner(slide, "Solution Overview", "Code modules mapped to product capabilities")
add_title(slide, "Unified Operations Workspace")
modules = [
    ("Dashboard", "KPI cards, alerts, activity feed"),
    ("Heatmap", "Need type layers + underserved zones"),
    ("Tasks", "Status lifecycle: Open -> Assigned -> In Progress -> Completed -> Verified"),
    ("Assignments", "Smart matching + manual override + reassignment logs"),
    ("Analytics", "Trend charts, category split, zone performance"),
    ("Volunteers", "Directory, availability, profile-level operations"),
]
start_y = Inches(1.9)
for i, (name, detail) in enumerate(modules):
    row = i // 2
    col = i % 2
    x = Inches(0.7) + col * Inches(6.35)
    y = start_y + row * Inches(1.35)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.9), Inches(1.1))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = ACCENT
    t = card.text_frame
    t.clear()
    p1 = t.paragraphs[0]
    p1.text = name
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    p2 = t.add_paragraph()
    p2.text = detail
    p2.font.size = Pt(11)
    p2.font.color.rgb = MUTED

# Slide 4: Architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_banner(slide, "System Architecture", "Built from current Flutter implementation")
add_title(slide, "End-to-End Data & Decision Flow")
box_w = Inches(2.1)
box_h = Inches(0.95)
x0 = Inches(0.45)
y0 = Inches(2.3)
labels = [
    "Field Reports\n& Coordinator Input",
    "Needs / Tasks\nData Models",
    "Smart Match Engine\n(skill 60% + proximity 30% + rating 10%)",
    "Assignment Engine\n(dialog + timeline)",
    "Analytics Layer\n(charts + zone table)",
    "Coordinator Actions\n(reassign, verify, escalate)",
]
for i, label in enumerate(labels):
    add_flow_box(slide, x0 + Inches(2.15 * i), y0, box_w, box_h, label)
    if i < len(labels) - 1:
        add_arrow(slide, x0 + Inches(2.15 * i) + box_w - Inches(0.04), y0 + Inches(0.38), Inches(0.22), Inches(0.16))

feedback = add_flow_box(slide, Inches(4.4), Inches(4.25), Inches(4.4), Inches(0.95), "Feedback Loop: completion + verification updates KPIs and zone coverage", TEAL)

# Slide 5: User Journey Flow
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_banner(slide, "Operational Flow", "Coordinator-first interaction journey")
add_title(slide, "From Need Detection to Verified Completion")
steps = [
    "Need Reported",
    "Need Prioritized",
    "Task Created",
    "Smart Match Suggested",
    "Volunteer Assigned",
    "Task Executed",
    "Coordinator Verifies",
    "Impact Logged",
]
x = Inches(0.55)
y = Inches(2.0)
for i, s in enumerate(steps):
    w = Inches(1.42)
    add_flow_box(slide, x + Inches(i * 1.58), y + (Inches(0.75) if i % 2 else Inches(0)), w, Inches(0.72), s, PRIMARY)
    if i < len(steps) - 1:
        add_arrow(slide, x + Inches(i * 1.58) + w - Inches(0.03), y + (Inches(0.28) if i % 2 == 0 else Inches(1.02)), Inches(0.2), Inches(0.14))

add_bullets(
    slide,
    Inches(0.7), Inches(4.7), Inches(12.0), Inches(1.7),
    [
        "Reassignment reasons are stored in timeline logs to improve transparency and accountability.",
        "Overdue and unassigned filters help coordinators triage quickly.",
        "Verified status closes the loop for downstream reporting and donor communication.",
    ],
)

# Slide 6: Heatmap + Geo Intelligence
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_banner(slide, "Geospatial Intelligence", "Need type overlays + underserved detection")
add_title(slide, "Heatmap Module")
add_bullets(
    slide,
    Inches(0.6), Inches(1.9), Inches(5.1), Inches(4.7),
    [
        "Need layers: food, medical, shelter, education.",
        "Volunteer overlay toggle for supply-demand visibility.",
        "Underserved zone circles when volunteers < 2 and intensity > 0.6.",
        "Time filtering: past 7 days vs past 30 days.",
        "Zone detail panel shows ward, intensity, and volunteer count.",
    ],
    title="Implemented Signals"
)
# pseudo map panel
map_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.95), Inches(1.9), Inches(6.85), Inches(4.85))
map_panel.fill.solid()
map_panel.fill.fore_color.rgb = WHITE
map_panel.line.color.rgb = RGBColor(0xC7, 0xD2, 0xE0)

# markers
coords = [
    (6.6, 2.7, RED), (7.9, 3.2, BLUE), (9.2, 2.6, YELLOW),
    (10.5, 3.3, GREEN), (11.3, 2.4, RED), (8.8, 4.4, BLUE),
    (10.0, 4.9, YELLOW), (7.1, 4.8, RED),
]
for cx, cy, c in coords:
    m = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(0.26), Inches(0.26))
    m.fill.solid()
    m.fill.fore_color.rgb = c
    m.line.fill.background()

vol_tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.2), Inches(5.7), Inches(2.0), Inches(0.48))
vol_tag.fill.solid()
vol_tag.fill.fore_color.rgb = RGBColor(0xDC, 0xFC, 0xE7)
vol_tag.line.fill.background()
vol_tag.text_frame.text = "4 Volunteers Active"
vol_tag.text_frame.paragraphs[0].font.size = Pt(10)

# Slide 7: Smart Matching
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_banner(slide, "Smart Assignment", "Ranking logic implemented in volunteer_model.dart")
add_title(slide, "Volunteer Recommendation Engine")
add_bullets(
    slide,
    Inches(0.6), Inches(1.8), Inches(6.1), Inches(3.8),
    [
        "Eligible set: only volunteers with `availability == available`.",
        "Score = (skill overlap * 0.6) + (distance score * 0.3) + (rating * 0.1).",
        "Distance scoring caps at 20 km to avoid long-distance assignments.",
        "UI supports smart-only view and full-volunteer manual override.",
        "Best match highlight plus % match explanation improves trust in AI suggestion.",
    ],
)

chart_data = CategoryChartData()
chart_data.categories = ["Riya", "Anjali", "Sneha", "Rohan"]
chart_data.add_series("Match %", (92, 88, 83, 78))
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(6.8), Inches(2.0), Inches(6.0), Inches(3.4),
    chart_data,
).chart
chart.has_legend = False
chart.value_axis.maximum_scale = 100
chart.value_axis.minimum_scale = 0
chart.value_axis.has_major_gridlines = True
chart.category_axis.tick_labels.font.size = Pt(11)
chart.value_axis.tick_labels.font.size = Pt(10)
chart.chart_title.has_text_frame = True
chart.chart_title.text_frame.text = "Top Smart-Match Candidates (Sample Task)"

# Slide 8: Analytics
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_banner(slide, "Analytics & Reporting", "Charts sourced from analytics_screen.dart mock datasets")
add_title(slide, "Performance Insights (This Month)")

line_data = CategoryChartData()
line_data.categories = [f"W{i}" for i in range(1, 13)]
line_data.add_series("Tasks Resolved", (12, 18, 14, 22, 19, 25, 21, 28, 24, 30, 27, 35))
line_chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE_MARKERS,
    Inches(0.6), Inches(1.9), Inches(7.2), Inches(3.1),
    line_data,
).chart
line_chart.has_legend = False
line_chart.value_axis.has_major_gridlines = True
line_chart.chart_title.has_text_frame = True
line_chart.chart_title.text_frame.text = "Task Resolution Trend"

pie_data = CategoryChartData()
pie_data.categories = ["Food", "Medical", "Shelter", "Education"]
pie_data.add_series("Needs", (42, 31, 19, 14))
pie = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE,
    Inches(8.1), Inches(1.9), Inches(4.7), Inches(3.1),
    pie_data,
).chart
pie.has_legend = True
pie.legend.position = XL_LEGEND_POSITION.RIGHT
pie.legend.include_in_layout = False
pie.chart_title.has_text_frame = True
pie.chart_title.text_frame.text = "Needs by Category"

bar_data = CategoryChartData()
bar_data.categories = [f"W{i}" for i in range(1, 13)]
bar_data.add_series("Active Volunteers", (20, 28, 24, 32, 27, 35, 30, 38, 33, 42, 38, 45))
bar = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.6), Inches(5.15), Inches(12.2), Inches(1.8),
    bar_data,
).chart
bar.has_legend = False
bar.chart_title.has_text_frame = True
bar.chart_title.text_frame.text = "Volunteer Activity Trend"

# Slide 9: Impact and KPIs
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_banner(slide, "Projected Impact", "How this system improves field outcomes")
add_title(slide, "Expected Operational Gains")

impact_data = CategoryChartData()
impact_data.categories = ["Response Time", "Coverage", "Match Accuracy", "Verification Completion"]
impact_data.add_series("Before", (100, 52, 61, 48))
impact_data.add_series("After", (62, 74, 86, 81))
impact_chart = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED,
    Inches(0.7), Inches(2.0), Inches(7.2), Inches(4.2),
    impact_data,
).chart
impact_chart.has_legend = True
impact_chart.legend.position = XL_LEGEND_POSITION.BOTTOM
impact_chart.value_axis.has_major_gridlines = True

add_kpi_card(slide, Inches(8.4), Inches(2.1), Inches(4.5), Inches(1.25), "Faster Assignment Cycle", "~38% reduction", GREEN)
add_kpi_card(slide, Inches(8.4), Inches(3.55), Inches(4.5), Inches(1.25), "Underserved Zone Reduction", "5 -> 2 zones", BLUE)
add_kpi_card(slide, Inches(8.4), Inches(5.0), Inches(4.5), Inches(1.25), "Volunteer Utilization Lift", "+22%", TEAL)

# Slide 10: Roadmap
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_banner(slide, "Execution Roadmap", "Prototype -> Pilot -> Scale")
add_title(slide, "Build & Rollout Plan")

milestones = [
    ("Phase 1 (Now)", "Frontend prototype complete\nDashboard, assignments, heatmap, analytics UI", GREEN),
    ("Phase 2", "Backend APIs + auth + report ingestion pipeline", BLUE),
    ("Phase 3", "Pilot in 2 districts with coordinator training", YELLOW),
    ("Phase 4", "Impact instrumentation + donor-grade reporting exports", TEAL),
]
for i, (m, d, c) in enumerate(milestones):
    y = Inches(1.9 + i * 1.28)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y + Inches(0.24), Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = c
    dot.line.fill.background()

    if i < len(milestones) - 1:
        v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.98), y + Inches(0.42), Inches(0.02), Inches(1.1))
        v.fill.solid()
        v.fill.fore_color.rgb = RGBColor(0xCF, 0xDA, 0xE6)
        v.line.fill.background()

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.25), y, Inches(11.5), Inches(1.02))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = c
    tf = card.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = m
    p1.font.size = Pt(15)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    p2 = tf.add_paragraph()
    p2.text = d
    p2.font.size = Pt(11)
    p2.font.color.rgb = MUTED

# Slide 11: Tech Stack & Governance
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_banner(slide, "Technology & Governance", "Scalable and accountable by design")
add_title(slide, "Stack, Data, and Reliability Considerations")
add_bullets(
    slide,
    Inches(0.6), Inches(1.9), Inches(6.3), Inches(4.7),
    [
        "Frontend: Flutter web admin (`dashboard`, `tasks`, `assignments`, `analytics`, `heatmap`).",
        "Visualization: `fl_chart` for KPI insights and trend charts.",
        "Geo layer: `flutter_map` + OpenStreetMap tiles for locality-level context.",
        "Data model supports full task lifecycle and assignment history logs.",
        "Next: role-based access, audit logs, secure backend sync, and CSV/PDF exports.",
    ],
)
add_flow_box(slide, Inches(7.1), Inches(2.1), Inches(5.5), Inches(1.0), "Security: coordinator roles + access control + audit history", BLUE)
add_flow_box(slide, Inches(7.1), Inches(3.45), Inches(5.5), Inches(1.0), "Reliability: offline-safe reports + retryable assignment actions", GREEN)
add_flow_box(slide, Inches(7.1), Inches(4.8), Inches(5.5), Inches(1.0), "Scalability: API-driven services + zone-level indexing", TEAL)

# Slide 12: Demo Script
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_banner(slide, "Prototype Demonstration", "Suggested live walkthrough sequence")
add_title(slide, "Demo Flow for Judges / Stakeholders")
add_bullets(
    slide,
    Inches(0.6), Inches(1.9), Inches(12.1), Inches(4.8),
    [
        "1. Open dashboard: show volunteers, active tasks, pending needs, and coverage KPIs.",
        "2. Switch to heatmap: highlight underserved circles and need-type toggles.",
        "3. Open assignments: use smart match for a task and explain score breakdown.",
        "4. Move task status through in-progress to verified and show timeline updates.",
        "5. End in analytics: show monthly trend lines, category mix, and zone performance.",
        "6. Close with impact story: faster response, better coverage, accountable execution.",
    ],
)

prs.save(OUT_FILE)
print(f"Created {OUT_FILE}")
